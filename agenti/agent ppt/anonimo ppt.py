from pptx import Presentation
from pathlib import Path
import zipfile
import json
import shutil
import xml.etree.ElementTree as ET
import entity as et
import output as ot


def extract_pptx_content(pptx_path: str, output_dir: str = "pptx_extracted"):
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)

    if output_dir.exists():
        shutil.rmtree(output_dir)

    output_dir.mkdir(parents=True, exist_ok=True)

    result = {
        "file_name": pptx_path.name,
        "slides": [],
        "metadata": {},
        "images": [],
        "links": []
    }

    # 1. Estrazione testo visibile dalle slide
    prs = Presentation(pptx_path)

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "slide_number": slide_index,
            "texts": [],
            "tables": [],
            "notes": []
        }

        for shape in slide.shapes:
            # Testo normale
            if hasattr(shape, "text") and shape.text.strip():
                slide_data["texts"].append(shape.text.strip())

            # Tabelle
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)

                slide_data["tables"].append(table_data)

        result["slides"].append(slide_data)

    # 2. Estrazione contenuto grezzo OpenXML
    raw_dir = output_dir / "raw_pptx"
    with zipfile.ZipFile(pptx_path, "r") as zip_ref:
        zip_ref.extractall(raw_dir)

    # 3. Estrazione metadata
    metadata_files = [
        raw_dir / "docProps" / "core.xml",
        raw_dir / "docProps" / "app.xml"
    ]

    for metadata_file in metadata_files:
        if metadata_file.exists():
            result["metadata"][metadata_file.name] = metadata_file.read_text(
                encoding="utf-8",
                errors="ignore"
            )

    # 4. Estrazione immagini
    media_dir = raw_dir / "ppt" / "media"

    if media_dir.exists():
        images_dir = output_dir / "images"
        images_dir.mkdir(exist_ok=True)

        for media_file in media_dir.iterdir():
            target_path = images_dir / media_file.name
            shutil.copy(media_file, target_path)

            result["images"].append({
                "file_name": media_file.name,
                "path": str(target_path)
            })

    # 5. Estrazione link embedded dai file .rels
    rels_files = list(raw_dir.rglob("*.rels"))

    for rels_file in rels_files:
        try:
            tree = ET.parse(rels_file)
            root = tree.getroot()

            for rel in root:
                target = rel.attrib.get("Target", "")

                if target.startswith("http") or target.startswith("mailto:"):
                    result["links"].append({
                        "source_file": str(rels_file.relative_to(raw_dir)),
                        "target": target
                    })

        except ET.ParseError:
            pass

    # 6. Salvataggio JSON riepilogativo
    json_path = output_dir / "pptx_content.json"

    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(result, f, indent=2, ensure_ascii=False)

    return result


# Estraggo contenuto ppt
data = extract_pptx_content("presentazione.pptx")
print(json.dumps(data, indent=2, ensure_ascii=False))

# Rileva dati sensibili
result = et.detect_sensitive_entities(
    input_json_path = "pptx_extracted/pptx_content.json",
    output_dir = "pptx_detection",
    language_model = "en_core_web_lg"
)

print(json.dumps(result, indent = 2, ensure_ascii = False))

# Genera l'output
output = ot.anonymize_presentation(
    input_pptx_path = "presentazione.pptx",
    anonymization_map_path = "pptx_detection/anonymization_map.json",
    output_pptx_path = "presentazione_anonimizzata.pptx",
    deep_openxml_pass = True
)

print(f"PowerPoint anonimizzato creato: {output}")