import json
import re
import zipfile
import shutil
from pathlib import Path
from pptx import Presentation


def load_anonymization_map(map_path: str):
    with open(map_path, "r", encoding="utf-8") as f:
        raw_map = json.load(f)

    replacements = {}

    for item in raw_map.values():
        original = item["original"]
        replacement = item["replacement"]
        replacements[original] = replacement

    return replacements


def replace_text_preserving_case(text: str, replacements: dict):
    if not text:
        return text

    # Prima sostituisce le entità più lunghe, per evitare collisioni
    sorted_items = sorted(
        replacements.items(),
        key=lambda x: len(x[0]),
        reverse=True
    )

    for original, replacement in sorted_items:
        pattern = re.compile(re.escape(original), re.IGNORECASE)
        text = pattern.sub(replacement, text)

    return text


def anonymize_text_frame(text_frame, replacements: dict):
    """
    Sostituisce testo dentro shape mantenendo il più possibile formattazione,
    paragrafi e run.
    """
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = replace_text_preserving_case(run.text, replacements)


def anonymize_slide_shapes(slide, replacements: dict):
    for shape in slide.shapes:
        # Testo normale
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            anonymize_text_frame(shape.text_frame, replacements)

        # Tabelle
        if hasattr(shape, "has_table") and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        anonymize_text_frame(cell.text_frame, replacements)

        # Gruppi di shape
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for grouped_shape in shape.shapes:
                if hasattr(grouped_shape, "has_text_frame") and grouped_shape.has_text_frame:
                    anonymize_text_frame(grouped_shape.text_frame, replacements)

                if hasattr(grouped_shape, "has_table") and grouped_shape.has_table:
                    for row in grouped_shape.table.rows:
                        for cell in row.cells:
                            anonymize_text_frame(cell.text_frame, replacements)


def anonymize_speaker_notes(slide, replacements: dict):
    """
    python-pptx supporta accesso alle note, ma non sempre in modo completo.
    Questo copre i casi standard.
    """
    if not slide.has_notes_slide:
        return

    notes_slide = slide.notes_slide

    for shape in notes_slide.shapes:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            anonymize_text_frame(shape.text_frame, replacements)


def anonymize_pptx_content(
    input_pptx_path: str,
    anonymization_map_path: str,
    output_pptx_path: str = "presentazione_anonimizzata.pptx"
):
    input_pptx_path = Path(input_pptx_path)
    output_pptx_path = Path(output_pptx_path)

    replacements = load_anonymization_map(anonymization_map_path)

    prs = Presentation(input_pptx_path)

    # Slide + tabelle + gruppi + note
    for slide in prs.slides:
        anonymize_slide_shapes(slide, replacements)
        anonymize_speaker_notes(slide, replacements)

    # Metadata accessibili da python-pptx
    core_props = prs.core_properties

    metadata_fields = [
        "author",
        "category",
        "comments",
        "content_status",
        "identifier",
        "keywords",
        "language",
        "last_modified_by",
        "subject",
        "title",
        "version"
    ]

    for field in metadata_fields:
        try:
            value = getattr(core_props, field)
            if isinstance(value, str):
                setattr(
                    core_props,
                    field,
                    replace_text_preserving_case(value, replacements)
                )
        except Exception:
            pass

    prs.save(output_pptx_path)

    return output_pptx_path


def anonymize_openxml_metadata_safe(
    pptx_path: str,
    replacements: dict,
    output_pptx_path: str
):
    pptx_path = Path(pptx_path)
    output_pptx_path = Path(output_pptx_path)

    temp_dir = output_pptx_path.parent / "_pptx_xml_tmp"

    if temp_dir.exists():
        shutil.rmtree(temp_dir)

    temp_dir.mkdir(parents=True, exist_ok=True)

    # Extract
    with zipfile.ZipFile(pptx_path, "r") as zip_ref:
        zip_ref.extractall(temp_dir)

    xml_files = list(temp_dir.rglob("*.xml"))
    xml_files += list(temp_dir.rglob("*.rels"))

    for xml_file in xml_files:
        try:
            tree = ET.parse(xml_file)
            root = tree.getroot()

            # modifica SOLO testo nodi
            for elem in root.iter():

                if elem.text:
                    elem.text = replace_text_preserving_case(
                        elem.text,
                        replacements
                    )

                if elem.tail:
                    elem.tail = replace_text_preserving_case(
                        elem.tail,
                        replacements
                    )

            tree.write(
                xml_file,
                encoding="utf-8",
                xml_declaration=True
            )

        except Exception:
            # alcuni XML possono essere particolari
            pass

    # rebuild pptx
    with zipfile.ZipFile(
        output_pptx_path,
        "w",
        zipfile.ZIP_DEFLATED
    ) as zip_out:

        for file_path in temp_dir.rglob("*"):

            if file_path.is_file():
                zip_out.write(
                    file_path,
                    file_path.relative_to(temp_dir)
                )

    shutil.rmtree(temp_dir)

    return output_pptx_path


def anonymize_openxml_metadata(
    pptx_path: str,
    replacements: dict,
    output_pptx_path: str
):
    """
    Secondo passaggio più profondo:
    apre il .pptx come zip e sostituisce entità residue nei file XML.
    Utile per metadata, relazioni, alt text, commenti e proprietà non coperte da python-pptx.
    """
    pptx_path = Path(pptx_path)
    output_pptx_path = Path(output_pptx_path)

    temp_dir = output_pptx_path.parent / "_pptx_openxml_tmp"

    if temp_dir.exists():
        shutil.rmtree(temp_dir)

    temp_dir.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(pptx_path, "r") as zip_ref:
        zip_ref.extractall(temp_dir)

    xml_like_extensions = {
        ".xml",
        ".rels"
    }

    for file_path in temp_dir.rglob("*"):
        if file_path.suffix.lower() not in xml_like_extensions:
            continue

        try:
            content = file_path.read_text(encoding="utf-8", errors="ignore")
            new_content = replace_text_preserving_case(content, replacements)

            if new_content != content:
                file_path.write_text(new_content, encoding="utf-8")

        except Exception:
            pass

    if output_pptx_path.exists():
        output_pptx_path.unlink()

    with zipfile.ZipFile(output_pptx_path, "w", zipfile.ZIP_DEFLATED) as zip_out:
        for file_path in temp_dir.rglob("*"):
            if file_path.is_file():
                zip_out.write(
                    file_path,
                    file_path.relative_to(temp_dir)
                )

    shutil.rmtree(temp_dir)

    return output_pptx_path


def anonymize_presentation(
    input_pptx_path: str,
    anonymization_map_path: str,
    output_pptx_path: str = "presentazione_anonimizzata.pptx",
    deep_openxml_pass: bool = True
):
    replacements = load_anonymization_map(anonymization_map_path)

    intermediate_path = Path(output_pptx_path).with_name(
        Path(output_pptx_path).stem + "_intermediate.pptx"
    )

    anonymize_pptx_content(
        input_pptx_path=input_pptx_path,
        anonymization_map_path=anonymization_map_path,
        output_pptx_path=intermediate_path
    )

    if deep_openxml_pass:
        anonymize_openxml_metadata_safe(
            pptx_path=intermediate_path,
            replacements=replacements,
            output_pptx_path=output_pptx_path
        )

        if intermediate_path.exists():
            intermediate_path.unlink()
    else:
        shutil.move(intermediate_path, output_pptx_path)

    return output_pptx_path